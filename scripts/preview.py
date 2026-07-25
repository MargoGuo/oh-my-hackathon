#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""preview.py — SVG → png 快速预览(模拟 pptx 效果)。

读 SVG → svg2pptx 转(内存)→ Pillow 渲染每页 png。改完 SVG 跑这个,
秒看效果(接近实际 pptx,含 margin=0/文本定位等转换逻辑),不用开 PowerPoint。

用法:
    S=${CLAUDE_PLUGIN_ROOT}/skills/oh-my-hackathon/scripts
    python "$S/preview.py" --input <svg目录或文件> --output preview/
    # 然后浏览器/图片查看器打开 preview/preview_NN.png

依赖: python-pptx + Pillow(缺则提示)。
非 Windows: 改下方 FONT_CANDIDATES 为系统字体路径。
"""
from __future__ import annotations
import argparse
import sys
from pathlib import Path

_HERE = Path(__file__).resolve().parent
if str(_HERE) not in sys.path:
    sys.path.insert(0, str(_HERE))

_PREFIX = "[oh-my-hackathon/preview]"

try:
    from pptx import Presentation
    from pptx.util import Emu
    from pptx.enum.text import PP_ALIGN
    from PIL import Image, ImageDraw, ImageFont
except ImportError as e:  # pragma: no cover
    print(f"{_PREFIX} 缺依赖: {e}\n  pip install python-pptx Pillow", file=sys.stderr)
    sys.exit(3)

import svg2pptx  # 同目录

# 字体候选(按优先级)。非 Windows 改这里。
_FONT_CANDIDATES = {
    "bold": [
        r"C:\Windows\Fonts\msyhbd.ttc",      # 微软雅黑 Bold(中文)
        "/System/Library/Fonts/PingFang.ttc",  # macOS
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",  # Linux
    ],
    "arial": [r"C:\Windows\Fonts\arialbd.ttf"],
    "cambria": [r"C:\Windows\Fonts\cambriab.ttf"],
    "consolas": [r"C:\Windows\Fonts\consola.ttf"],
}


def _ep(v):
    return (v or 0) // 9525


def _fc(f):
    try:
        if f.type == 1:
            return "#" + str(f.fore_color.rgb)
    except Exception:
        pass
    return None


def _font(sz_pt, name):
    """按字体名选候选,回退到默认。"""
    pools = []
    if name and "Arial" in name:
        pools.append(_FONT_CANDIDATES["arial"])
    if name and "Cambria" in name:
        pools.append(_FONT_CANDIDATES["cambria"])
    if name and "Consolas" in name:
        pools.append(_FONT_CANDIDATES["consolas"])
    pools.append(_FONT_CANDIDATES["bold"])
    sz_px = max(int(sz_pt * 1.33), 8)
    for pool in pools:
        for fp in pool:
            try:
                return ImageFont.truetype(fp, sz_px)
            except Exception:
                continue
    return ImageFont.load_default()


def render_slide(slide) -> Image.Image:
    """渲染一页 slide → PIL Image(1280x720)。模拟 pptx 实际效果。"""
    img = Image.new("RGB", (1280, 720), "#FFFFFF")
    d = ImageDraw.Draw(img)
    for sh in slide.shapes:
        x, y, w, h = _ep(sh.left), _ep(sh.top), _ep(sh.width), _ep(sh.height)
        t = type(sh).__name__
        if t in ("Shape", "AutoShape"):
            c = _fc(sh.fill)
            outline = "black"
            try:
                if sh.line.color.type:
                    outline = "#" + str(sh.line.color.rgb)
            except Exception:
                pass
            d.rectangle([x, y, x + w, y + h], fill=c, outline=outline)
        elif t == "Picture":
            d.rectangle([x, y, x + w, y + h], outline="blue")
            d.text((x + 8, y + 8), "[图]", fill="blue")
        elif t == "Connector":
            try:
                d.line([x, y, x + w, y + h], fill="black", width=3)
            except Exception:
                pass
        if sh.has_text_frame:
            for pa in sh.text_frame.paragraphs:
                txt = "".join(r.text for r in pa.runs)
                if not txt.strip():
                    continue
                r0 = pa.runs[0]
                fs = r0.font.size.pt if r0.font.size else 18
                try:
                    col = "#" + str(r0.font.color.rgb)
                except Exception:
                    col = "#1F2937"
                f = _font(fs, r0.font.name or "")
                al = pa.alignment
                if al == PP_ALIGN.CENTER:
                    bb = d.textbbox((0, 0), txt, font=f)
                    tx = x + (w - (bb[2] - bb[0])) // 2
                elif al == PP_ALIGN.RIGHT:
                    bb = d.textbbox((0, 0), txt, font=f)
                    tx = x + w - (bb[2] - bb[0])
                else:
                    tx = x
                d.text((tx, y), txt, font=f, fill=col)
    return img


def main() -> int:
    p = argparse.ArgumentParser(prog="preview.py", description="SVG→png 预览(模拟 pptx)")
    p.add_argument("--input", required=True, help="SVG 文件或目录")
    p.add_argument("--output", required=True, help="png 输出目录")
    args = p.parse_args()

    src = Path(args.input)
    out = Path(args.output)
    out.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Emu(1280 * 9525)
    prs.slide_height = Emu(720 * 9525)

    svgs = [src] if src.is_file() else sorted(src.glob("*.svg"))
    if not svgs:
        print(f"{_PREFIX} 无 SVG: {src}", file=sys.stderr)
        return 2

    for svg in svgs:
        svg2pptx.convert_svg_file(prs, svg)

    for i, slide in enumerate(prs.slides, 1):
        render_slide(slide).save(out / f"preview_{i:02d}.png")
    print(f"{_PREFIX} ✅ {len(prs.slides)} 页 → {out}/")
    return 0


if __name__ == "__main__":
    try:
        sys.exit(main())
    except SystemExit:
        raise
    except Exception as exc:
        print(f"{_PREFIX} 未预期错误: {exc}", file=sys.stderr)
        sys.exit(4)

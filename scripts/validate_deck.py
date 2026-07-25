#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""validate_deck.py — 验证生成的 PPTX(借鉴 guizang validate + overflowFix 思路)。

检查每页:形状超界 / 字号低于投影下限 / 灰字对比不足。每条给修复建议。
用法:
    python scripts/validate_deck.py <path.pptx>
    # 返回 0=通过, 1=有错误(超界/字号), 警告不计错
"""
from __future__ import annotations
import sys
from pathlib import Path

_PREFIX = "[validate]"

try:
    from pptx import Presentation
except ImportError:  # pragma: no cover
    print(f"{_PREFIX} 缺 python-pptx。pip install python-pptx", file=sys.stderr)
    sys.exit(3)

# 画布(px)。SVG viewBox 1280x720 约定。
CANVAS_W, CANVAS_H = 1280, 720
# 字号下限(pt,投影清晰)。导航/标签 14,正文 ≥17。
MIN_FONT_PX = 14
# 对比不足的灰(浅底 + 小字号 → 投影模糊)。改 #0A0A0A 或 #1F2937。
GRAY_LOW_CONTRAST = {"44403C", "6B7280"}


def _overflow_fix(px_over):
    """借鉴 guizang overflowFix:按超界 px 给分级修复建议。"""
    n = round(px_over)
    if n <= 40:
        return f"超 {n}px:上移内容 或 收紧一个 gap/padding 20-40px,别删内容"
    if n <= 90:
        return f"超 {n}px:压缩局部 gap/padding + 减一个块高,避免删文案"
    if n <= 160:
        return f"超 {n}px:略缩标题 或 压一段,再考虑删"
    return f"超 {n}px:换高容量版式 或 有意删/合并内容"


def validate(pptx_path):
    """返回 (errors, warnings)。errors=必改(超界/字号),warnings=建议(灰字)。"""
    p = Presentation(pptx_path)
    errors, warnings = [], []
    for i, slide in enumerate(p.slides, 1):
        for sh in slide.shapes:
            l = (sh.left or 0) // 9525
            t = (sh.top or 0) // 9525
            w = (sh.width or 0) // 9525
            h = (sh.height or 0) // 9525
            # 超界
            if l + w > CANVAS_W + 2:
                errors.append(f"p{i}: 右侧超界({l+w}>{CANVAS_W}) — {_overflow_fix(l+w-CANVAS_W)}")
            if t + h > CANVAS_H + 2:
                errors.append(f"p{i}: 底部超界({t+h}>{CANVAS_H}) — {_overflow_fix(t+h-CANVAS_H)}")
            if l < -2:
                errors.append(f"p{i}: 左侧超界({l}) — 检查 g transform translate")
            # 文本
            if sh.has_text_frame:
                for para in sh.text_frame.paragraphs:
                    for run in para.runs:
                        fs_pt = run.font.size.pt if run.font.size else None
                        fs_px = (fs_pt * 1.33) if fs_pt else None
                        txt = (run.text or "").strip()
                        if not txt:
                            continue
                        if fs_px and fs_px < MIN_FONT_PX:
                            errors.append(f"p{i}: 字号 {fs_pt}pt(<{MIN_FONT_PX}px 下限) — '{txt[:15]}'")
                        try:
                            col = str(run.font.color.rgb)
                            if col in GRAY_LOW_CONTRAST and fs_px and fs_px < 16:
                                warnings.append(f"p{i}: 灰字 #{col} + 小字({fs_pt}pt) 对比不足 — 改 #0A0A0A")
                        except Exception:
                            pass
    return errors, warnings


def main() -> int:
    if len(sys.argv) < 2:
        print(f"用法: python scripts/validate_deck.py <pptx>", file=sys.stderr)
        return 2
    pptx = sys.argv[1]
    if not Path(pptx).exists():
        print(f"{_PREFIX} 文件不存在: {pptx}", file=sys.stderr)
        return 2
    errors, warnings = validate(pptx)
    print(f"{_PREFIX} {pptx}: {len(errors)} 错误, {len(warnings)} 警告")
    if errors:
        print("\n❌ 错误(必改):")
        for e in errors:
            print(f"  - {e}")
    if warnings:
        print("\n⚠️ 警告(建议):")
        for w in warnings:
            print(f"  - {w}")
    if not errors and not warnings:
        print(f"{_PREFIX} ✅ 全部通过")
    return 1 if errors else 0


if __name__ == "__main__":
    try:
        sys.exit(main())
    except SystemExit:
        raise
    except Exception as exc:
        print(f"{_PREFIX} 未预期错误: {exc}", file=sys.stderr)
        sys.exit(4)

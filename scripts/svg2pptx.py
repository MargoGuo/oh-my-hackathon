#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
svg2pptx.py — 精简 SVG → 可编辑 PPTX 转换器(oh-my-hackathon 自包含)。

设计借鉴 ppt-master(hugohe3/ppt-master, MIT)的核心思路:
  - SVG 作为中间格式: 每页一个 <svg viewBox="0 0 1280 720"> (16:9)
  - 每个 SVG 元素映射成一个 PowerPoint 原生可编辑形状(不栅格化成图片)
  - 圆角矩形保留为 roundRect(可继续调圆角), 不转成 path
  - CSS px → EMU 恒定映射: 1px = 9525 EMU (96px/in, 914400 EMU/in)

但本实现不复制 ppt-master 的 3 万行引擎, 而是用 python-pptx 高层 API
重新实现一个 ~400 行的精简版, 只覆盖黑客松路演 PPT 的必要能力:
  <rect>/<rect rx> → 矩形/圆角矩形
  <text>/<tspan>   → 文本框(可继续改字)
  <circle>/<ellipse> → 椭圆
  <line>           → 连接线
  <image href>     → 图片
  <g transform>    → 递归 + 平移

不支持的 SVG 特性(path 渐变 mask foreignObject 等)记 warning 跳过,
不中断。MVP 不内联图标库(<use data-icon> 记 warning)。
"""

from __future__ import annotations
import argparse
import base64
import re
import sys
from pathlib import Path
from xml.etree import ElementTree as ET

# ---------- 依赖探测(参考 study-bilibili-video 的 discover-not-assume 模式) ----------
try:
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
    PPTX_AVAILABLE = True
except ImportError:  # pragma: no cover
    PPTX_AVAILABLE = False

_HERE = Path(__file__).resolve().parent
_PREFIX = "[oh-my-hackathon/svg2pptx]"

# Windows 控制台默认 cp936, emoji 会 UnicodeEncodeError 导致误退出码 → reconfigure UTF-8
for _stream in (sys.stdout, sys.stderr):
    try:
        _stream.reconfigure(encoding="utf-8")
    except Exception:
        pass

EMU_PER_PX = 9525  # CSS px → EMU (标准: 96px=1in, 1in=914400 EMU)
SVG_W, SVG_H = 1280, 720  # 16:9 viewBox 约定
SVG_NS = "{http://www.w3.org/2000/svg}"
XLINK_HREF = "{http://www.w3.org/1999/xlink}href"

# CSS 命名色子集
_NAMED = {
    "black": "000000", "white": "FFFFFF", "red": "FF0000", "green": "008000",
    "blue": "0000FF", "yellow": "FFFF00", "orange": "FFA500", "purple": "800080",
    "gray": "808080", "grey": "808080", "transparent": "FFFFFF",
}


def _die(code: int, msg: str) -> None:
    print(f"{_PREFIX} {msg}", file=sys.stderr)
    sys.exit(code)


def _warn(msg: str) -> None:
    print(f"{_PREFIX} ⚠️ {msg}", file=sys.stderr)


def px(v: float) -> Emu:
    """SVG px → PPT EMU。"""
    return Emu(int(round(float(v) * EMU_PER_PX)))


def parse_color(s):
    """解析 fill/stroke 颜色 → RGBColor 或 None。支持 #rgb/#rrggbb/rgb()/rgba()/命名。"""
    if not s or s.strip().lower() in ("none", "transparent", "currentcolor"):
        return None
    s = s.strip()
    hexm = re.fullmatch(r"#([0-9a-fA-F]{3}|[0-9a-fA-F]{6})", s)
    if hexm:
        h = hexm.group(1)
        if len(h) == 3:
            h = "".join(c * 2 for c in h)
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    rgbm = re.fullmatch(r"rgba?\(([^)]+)\)", s)
    if rgbm:
        parts = [p.strip() for p in rgbm.group(1).split(",")][:3]
        vals = []
        for p in parts:
            if p.endswith("%"):
                vals.append(int(float(p[:-1]) * 255 / 100))
            else:
                vals.append(int(float(p)))
        return RGBColor(*vals[:3])
    if s.lower() in _NAMED:
        h = _NAMED[s.lower()]
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    return None


def parse_transform(s: str):
    """解析 transform → (dx, dy, sx, sy)。仅支持 translate/scale(rotate 忽略)。"""
    dx = dy = 0.0
    sx = sy = 1.0
    if not s:
        return dx, dy, sx, sy
    for m in re.finditer(r"(translate|scale)\(([^)]+)\)", s):
        fn, raw = m.group(1), m.group(2).replace(",", " ").split()
        vals = [float(a) for a in raw]
        if fn == "translate":
            dx += vals[0]
            dy += vals[1] if len(vals) > 1 else 0.0
        elif fn == "scale":
            sx *= vals[0]
            sy *= vals[1] if len(vals) > 1 else vals[0]
    return dx, dy, sx, sy


def _text_width(text: str, fs_px: float) -> float:
    """估算文字像素宽度,用于 text-anchor 居中:CJK/全角≈1.0em,窄字符≈0.3em,其他≈0.58em。"""
    w = 0.0
    for ch in text:
        o = ord(ch)
        if o > 0x2E80:  # CJK 统一汉字及以远(含全角标点)
            w += fs_px * 1.0
        elif ch in "iIl.,;:' `":
            w += fs_px * 0.3
        else:
            w += fs_px * 0.58
    return w


def style_of(elem) -> dict:
    """合并 element attribute + style 属性 → 样式 dict。attribute 优先级低(可被 style 覆盖)。"""
    d = {}
    style = elem.get("style")
    if style:
        for kv in style.split(";"):
            if ":" in kv:
                k, v = kv.split(":", 1)
                d[k.strip()] = v.strip()
    for k in ("fill", "stroke", "stroke-width", "font-size", "font-family",
              "font-weight", "text-anchor", "fill-opacity", "opacity",
              "rx", "ry", "transform"):
        v = elem.get(k)
        if v is not None and k not in d:
            d[k] = v
    return d


def _opacity_to_alpha(s):
    """opacity/fill-opacity (0~1) → 0~100000 (pptx alpha)。None → None。"""
    if s is None:
        return None
    try:
        return max(0, min(100000, int(round(float(s) * 100000))))
    except (TypeError, ValueError):
        return None


def _apply_alpha(shape, alpha):
    """给形状填充加透明度(alpha 为 None 跳过)。"""
    if alpha is None or alpha >= 100000:
        return
    try:
        from pptx.oxml.ns import qn
        srgb = shape.fill._fill.find(qn("a:srgbClr"))
        if srgb is not None:
            a = srgb.makeelement(qn("a:alpha"), {"val": str(alpha)})
            srgb.append(a)
    except Exception:
        pass


# ---------- 形状渲染 ----------
def add_rect(slide, x, y, w, h, st, tdx, tdy):
    rx = float(st.get("rx") or 0) or float(st.get("ry") or 0)
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rx > 0 else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, px(x + tdx), px(y + tdy), px(w), px(h))
    fill = parse_color(st.get("fill"))
    stroke = parse_color(st.get("stroke"))
    if fill is None:  # 默认黑填充太丑, 无 fill 时改成无填充
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        _apply_alpha(shp, _opacity_to_alpha(st.get("fill-opacity") or st.get("opacity")))
    if stroke is None or st.get("stroke-width") == "0":
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
        sw = st.get("stroke-width", "1")
        try:
            shp.line.width = Pt(float(sw) * 0.75)
        except (TypeError, ValueError):
            shp.line.width = Pt(0.75)
    shp.shadow.inherit = False  # 去掉默认阴影(投影更干净)
    # 圆角半径: pptx 的 adj 取 0~0.5(相对短边比例)
    if rx > 0:
        try:
            shp.adjustments[0] = min(0.5, rx / min(w, h))
        except Exception:
            pass
    return shp


def add_ellipse(slide, x, y, w, h, st, tdx, tdy):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, px(x + tdx), px(y + tdy), px(w), px(h))
    fill = parse_color(st.get("fill"))
    stroke = parse_color(st.get("stroke"))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if stroke is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = stroke
    shp.shadow.inherit = False
    return shp


def add_line(slide, x1, y1, x2, y2, st, tdx, tdy):
    from pptx.enum.shapes import MSO_CONNECTOR
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      px(x1 + tdx), px(y1 + tdy),
                                      px(x2 + tdx), px(y2 + tdy))
    stroke = parse_color(st.get("stroke"))
    if stroke is None:
        conn.line.color.rgb = RGBColor(0, 0, 0)
    else:
        conn.line.color.rgb = stroke
    try:
        conn.line.width = Pt(float(st.get("stroke-width", "1")) * 0.75)
    except (TypeError, ValueError):
        conn.line.width = Pt(0.75)
    return conn


def _set_run_font(run, st):
    """把 SVG font-* 样式写到 pptx run。"""
    fs = st.get("font-size")
    if fs:
        try:
            run.font.size = Pt(float(re.sub(r"[^0-9.]", "", fs)) * 0.75)
        except (TypeError, ValueError):
            pass
    color = parse_color(st.get("fill"))
    if color is not None:
        run.font.color.rgb = color
    fw = st.get("font-weight")
    if fw and fw != "normal":
        run.font.bold = True
    fam = st.get("font-family")
    if fam:
        run.font.name = fam.split(",")[0].strip().strip("'\"")


def add_text(slide, elem, st, tdx, tdy):
    """SVG <text> → pptx textbox。<tspan> 当多 run/多行处理。"""
    x = float(elem.get("x", 0))
    y = float(elem.get("y", 0))
    anchor = st.get("text-anchor", "start")
    # 估算字号用于垂直定位(text 基线 → 框顶部)
    fs_px = 16.0
    fsm = re.search(r"([0-9.]+)", st.get("font-size", "") or "")
    if fsm:
        fs_px = float(fsm.group(1))
    box_w = float(st.get("_width", 400))  # 调用方可注入文字框宽
    left = x + tdx - (box_w / 2 if anchor == "middle" else (0 if anchor == "start" else box_w))
    top = y + tdy - fs_px * 0.82  # 基线 → 框顶(ascent≈0.82 字号, 比满字号更准)
    tb = slide.shapes.add_textbox(px(left), px(top), px(box_w), px(fs_px * 1.35))
    tf = tb.text_frame
    tf.word_wrap = False  # 不换行: 避免估算宽度偏小导致长文字断成多行打乱布局
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0  # 去 textbox 默认内边距(左右 0.1"), 文字从框边界精确开始
    # 收集 text 直接内容 + tspan
    para = tf.paragraphs[0]
    para.alignment = {"start": PP_ALIGN.LEFT, "middle": PP_ALIGN.CENTER,
                      "end": PP_ALIGN.RIGHT}.get(anchor, PP_ALIGN.LEFT)
    tspans = list(elem.findall(f"{SVG_NS}tspan"))
    if tspans:
        first = True
        for ts in tspans:
            tst = style_of(ts)
            # _set_run_font 需要合并 text 级样式
            merged = {**st, **{k: v for k, v in tst.items() if k in st or k in
                               ("fill", "font-size", "font-weight", "font-family")}}
            txt = (ts.text or "")
            if first:
                run = para.add_run()
                first = False
            else:
                run = para.add_run()
            run.text = txt
            _set_run_font(run, merged)
    else:
        run = para.add_run()
        run.text = (elem.text or "").strip()
        _set_run_font(run, st)
    return tb


def add_image(slide, elem, st, tdx, tdy, base_dir: Path):
    """<image href> → pptx picture。支持 data:base64 和相对/绝对路径。"""
    href = elem.get(XLINK_HREF) or elem.get("href")
    if not href:
        return None
    x = float(elem.get("x", 0)); y = float(elem.get("y", 0))
    w = float(elem.get("width", 0)); h = float(elem.get("height", 0))
    if href.startswith("data:"):
        m = re.search(r"base64,(.+)", href)
        if not m:
            return None
        import tempfile, os
        data = base64.b64decode(m.group(1))
        suffix = ".png" if "png" in href[:30] else ".jpg"
        tmp = tempfile.NamedTemporaryFile(suffix=suffix, delete=False)
        tmp.write(data); tmp.close()
        pic = slide.shapes.add_picture(tmp.name, px(x + tdx), px(y + tdy),
                                       px(w) if w else None, px(h) if h else None)
        os.unlink(tmp.name)
        return pic
    # 文件路径
    p = Path(href)
    if not p.is_absolute():
        p = base_dir / p
    if not p.exists():
        _warn(f"image not found: {href}")
        return None
    return slide.shapes.add_picture(str(p), px(x + tdx), px(y + tdy),
                                    px(w) if w else None, px(h) if h else None)


# ---------- 主遍历 ----------
def render_elem(slide, elem, base_dir: Path, tdx=0.0, tdy=0.0):
    """递归渲染一个 SVG 元素。g 的 transform 累加到子元素偏移。"""
    tag = elem.tag.replace(SVG_NS, "")
    st = style_of(elem)
    # 本元素的 transform
    ndx, ndy, _sx, _sy = parse_transform(st.get("transform"))
    cdx, cdy = tdx + ndx, tdy + ndy

    if tag == "g":
        for child in elem:
            render_elem(slide, child, base_dir, cdx, cdy)
        return
    if tag == "rect":
        try:
            add_rect(slide, float(elem.get("x", 0)), float(elem.get("y", 0)),
                     float(elem.get("width", 0)), float(elem.get("height", 0)),
                     st, cdx, cdy)
        except (TypeError, ValueError) as e:
            _warn(f"rect skipped: {e}")
    elif tag in ("circle", "ellipse"):
        if tag == "circle":
            cx = float(elem.get("cx", 0)); cy = float(elem.get("cy", 0))
            r = float(elem.get("r", 0))
            add_ellipse(slide, cx - r, cy - r, 2 * r, 2 * r, st, cdx, cdy)
        else:
            cx = float(elem.get("cx", 0)); cy = float(elem.get("cy", 0))
            rx_ = float(elem.get("rx", 0)); ry_ = float(elem.get("ry", 0))
            add_ellipse(slide, cx - rx_, cy - ry_, 2 * rx_, 2 * ry_, st, cdx, cdy)
    elif tag == "line":
        add_line(slide, float(elem.get("x1", 0)), float(elem.get("y1", 0)),
                 float(elem.get("x2", 0)), float(elem.get("y2", 0)), st, cdx, cdy)
    elif tag == "text":
        # 估算文字实际宽度(中英文区分),用于 text-anchor 居中定位
        fsm = re.search(r"([0-9.]+)", st.get("font-size", "") or "")
        fs = float(fsm.group(1)) if fsm else 16.0
        parts = [elem.text or ""] + [(t.text or "") for t in elem.findall(f"{SVG_NS}tspan")]
        st["_width"] = max(_text_width("".join(parts), fs), 60)
        add_text(slide, elem, st, cdx, cdy)
    elif tag == "image":
        try:
            add_image(slide, elem, st, cdx, cdy, base_dir)
        except Exception as e:
            _warn(f"image skipped: {e}")
    elif tag == "use":
        _warn(f"<use data-icon> 不内联(无图标库): 跳过 {elem.get('href') or elem.get(XLINK_HREF)}")
    elif tag == "path":
        _warn("<path> 不支持(精简版): 跳过")  # TODO: 可选 custGeom
    elif tag in ("svg", "defs", "style", "title", "desc", "metadata"):
        if tag == "svg":
            for child in elem:
                render_elem(slide, child, base_dir, cdx, cdy)
    else:
        _warn(f"未知元素跳过: <{tag}>")


def convert_svg_file(prs, svg_path: Path):
    """把一个 SVG 文件渲染成新 slide。文件名 NN_xxx.svg 决定页序。"""
    try:
        tree = ET.parse(svg_path)
    except ET.ParseError as e:
        _warn(f"解析失败 {svg_path.name}: {e}")
        return
    root = tree.getroot()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 6 = blank
    render_elem(slide, root, svg_path.parent)


def convert_dir(svg_dir: Path, out_pptx: Path):
    if not PPTX_AVAILABLE:
        _die(3, "python-pptx 未安装。pip install python-pptx 后重试。")
    if not svg_dir.is_dir():
        _die(2, f"SVG 目录不存在: {svg_dir}")
    svgs = sorted(svg_dir.glob("*.svg"))
    if not svgs:
        _die(2, f"目录里没有 .svg: {svg_dir}")
    prs = Presentation()
    # 16:9 slide (13.333in x 7.5in = 12192000 x 6858000 EMU), 对齐 1280x720 viewBox
    prs.slide_width = Emu(1280 * EMU_PER_PX)
    prs.slide_height = Emu(720 * EMU_PER_PX)
    for svg in svgs:
        convert_svg_file(prs, svg)
    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_pptx))
    print(f"{_PREFIX} ✅ {len(svgs)} 页 → {out_pptx}")


def main() -> int:
    if not PPTX_AVAILABLE:
        print(f"{_PREFIX} ⚠️ python-pptx 未装, 仅打印诊断", file=sys.stderr)
    p = argparse.ArgumentParser(prog="svg2pptx.py",
                                description="精简 SVG→可编辑 PPTX(借鉴 ppt-master 思路)")
    p.add_argument("--input", required=True, help="SVG 文件或目录(目录按文件名排序)")
    p.add_argument("--output", required=True, help="输出 .pptx 路径")
    args = p.parse_args()

    src = Path(args.input)
    if src.is_file():
        if not PPTX_AVAILABLE:
            _die(3, "python-pptx 未安装")
        prs = Presentation()
        prs.slide_width = Emu(1280 * EMU_PER_PX)
        prs.slide_height = Emu(720 * EMU_PER_PX)
        convert_svg_file(prs, src)
        Path(args.output).parent.mkdir(parents=True, exist_ok=True)
        prs.save(args.output)
        print(f"{_PREFIX} ✅ 1 页 → {args.output}")
        return 0
    convert_dir(src, Path(args.output))
    return 0


if __name__ == "__main__":
    try:
        sys.exit(main())
    except SystemExit:
        raise
    except Exception as exc:
        _die(4, f"未预期错误: {exc}")
